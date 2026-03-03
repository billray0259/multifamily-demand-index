"""
Create Demand Index App Presentation for Quannah Partners.

Generates a 10-slide widescreen PowerPoint that explains what the
Multifamily Demand Index tool is, how to use it, the methodology
behind it, its limitations, validation results, and how it supports
expansion into new value-add multifamily markets.

Styling mirrors the Streamlit app and Excel export palette.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# Directory containing this script (for resolving figure paths)
_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

# ══════════════════════════════════════════════════════════════════════════════
# APP-CONSISTENT PALETTE  (from .streamlit/config.toml & export.py)
# ══════════════════════════════════════════════════════════════════════════════
PRIMARY     = RGBColor(0x1F, 0x4E, 0x79)   # header fill / primary buttons
ACCENT_DARK = RGBColor(0x1F, 0x38, 0x64)   # darker accent for formulas
SLIDE_BG    = RGBColor(0xFF, 0xFF, 0xFF)   # white — app background
SECONDARY_BG= RGBColor(0xF0, 0xF4, 0xF8)   # light blue-gray — sidebar / cards
TEXT_COLOR  = RGBColor(0x1A, 0x1A, 0x2E)   # main body text
NOTE_GRAY   = RGBColor(0x59, 0x59, 0x59)   # footnotes / muted text
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_FILL = RGBColor(0xD6, 0xE4, 0xF0)   # light-blue accent cells
ALT_ROW     = RGBColor(0xF2, 0xF7, 0xFC)   # alternating row tint
HIGH_GREEN  = RGBColor(0xC6, 0xEF, 0xCE)   # High Demand tier
MOD_YELLOW  = RGBColor(0xFF, 0xEB, 0x9C)   # Moderate Demand tier
LOW_RED     = RGBColor(0xFF, 0xC7, 0xCE)   # Low Demand tier
CHART_GREEN = RGBColor(0x4A, 0xDE, 0x80)   # chart bar green
CHART_RED   = RGBColor(0xF8, 0x71, 0x71)   # chart bar red
BORDER_GRAY = RGBColor(0xD9, 0xD9, 0xD9)   # thin borders

# Typography — Roboto for a prestigious, investment-grade look
FONT_HEAD = "Roboto"   # slide titles, section headers, bold labels
FONT_BODY = "Roboto"   # body text, bullets, callouts, table cells


# ══════════════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ══════════════════════════════════════════════════════════════════════════════


def _apply_bullet(p, char="•"):
    """Apply a real PowerPoint bullet to paragraph p.

    Correctly injects <a:buFont> + <a:buChar> into pPr in OOXML schema order,
    inserting them before <a:defRPr> so the element sequence is valid.
    Sets a hanging indent so wrapped lines align under the text, not the bullet.
    """
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    pPr = p._p.get_or_add_pPr()
    # Hanging indent: ~0.25 in in EMUs
    pPr.set("marL", "320040")
    pPr.set("indent", "-320040")

    # Remove any pre-existing bullet-related elements
    for tag in ("a:buNone", "a:buAutoNum", "a:buFont", "a:buChar"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)

    # Build buFont and buChar elements
    buFont = etree.Element(f"{{{NS_A}}}buFont")
    buFont.set("typeface", "Arial")
    buChar = etree.Element(f"{{{NS_A}}}buChar")
    buChar.set("char", char)

    # Insert before <a:defRPr> to respect OOXML child-element ordering
    defRPr = pPr.find(qn("a:defRPr"))
    if defRPr is not None:
        idx = list(pPr).index(defRPr)
        pPr.insert(idx, buChar)
        pPr.insert(idx, buFont)
    else:
        pPr.append(buFont)
        pPr.append(buChar)

def add_title_slide(prs, title, subtitle=None):
    """Navy-primary title slide matching the app header aesthetic."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Primary background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    # Thin accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    0, Inches(3.6),
                                    prs.slide_width, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_FILL
    accent.line.fill.background()

    # Title text
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.4),
                                  Inches(12.33), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = FONT_HEAD
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(3.8),
                                      Inches(12.33), Inches(1.0))
        tf = sb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.italic = True
        p.font.name = FONT_BODY
        p.font.color.rgb = ACCENT_FILL
        p.alignment = PP_ALIGN.CENTER

    # Dark footer band
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    0, Inches(6.8),
                                    prs.slide_width, Inches(0.7))
    footer.fill.solid()
    footer.fill.fore_color.rgb = ACCENT_DARK
    footer.line.fill.background()

    return slide


def add_content_slide(prs, title, subtitle=None):
    """Content slide with primary header bar and light body."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White body
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLIDE_BG
    bg.line.fill.background()

    # Primary header bar
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, 0, prs.slide_width, Inches(1.1))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = PRIMARY
    hdr.line.fill.background()

    # Accent underline
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, Inches(1.1),
                                 prs.slide_width, Inches(0.04))
    acc.fill.solid()
    acc.fill.fore_color.rgb = ACCENT_FILL
    acc.line.fill.background()

    # Title text
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15),
                                  Inches(12.5), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = FONT_HEAD
    p.font.color.rgb = WHITE

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.4), Inches(0.65),
                                      Inches(12.5), Inches(0.4))
        tf = sb.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.name = FONT_BODY
        p.font.color.rgb = ACCENT_FILL

    # Footer band
    fb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                0, Inches(7.05),
                                prs.slide_width, Inches(0.45))
    fb.fill.solid()
    fb.fill.fore_color.rgb = SECONDARY_BG
    fb.line.fill.background()

    ft = slide.shapes.add_textbox(Inches(0.4), Inches(7.1),
                                  Inches(8), Inches(0.3))
    p = ft.text_frame.paragraphs[0]
    p.text = "Quannah Partners  ·  Confidential"
    p.font.size = Pt(9)
    p.font.name = FONT_BODY
    p.font.color.rgb = NOTE_GRAY
    p.font.italic = True

    sn = slide.shapes.add_textbox(Inches(10), Inches(7.1),
                                  Inches(3), Inches(0.3))
    p = sn.text_frame.paragraphs[0]
    p.text = f"Slide {len(prs.slides)}"
    p.font.size = Pt(9)
    p.font.name = FONT_BODY
    p.font.color.rgb = NOTE_GRAY
    p.alignment = PP_ALIGN.RIGHT

    return slide


def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=TEXT_COLOR,
                 align=PP_ALIGN.LEFT, italic=False, font_name=None):
    """General-purpose text box."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name or (FONT_HEAD if bold else FONT_BODY)
    p.alignment = align
    return box


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=13, color=TEXT_COLOR, spacing=6):
    """Bulleted list using real PowerPoint bullet formatting."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True

    for i, raw_item in enumerate(items):
        # Strip any legacy text-character bullet prefixes
        item = raw_item.lstrip()
        for prefix in ("• ", "•", "– ", "–", "- "):
            if item.startswith(prefix):
                item = item[len(prefix):]
                break

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.name = FONT_BODY
        p.font.color.rgb = color
        p.space_before = Pt(spacing)
        _apply_bullet(p)
    return box


def add_callout_box(slide, text, left, top, width, height,
                    bg_color=ACCENT_FILL, text_color=ACCENT_DARK,
                    font_size=12, bold=False):
    """Rounded-rectangle callout matching the app's accent palette."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.name = FONT_BODY
    tf.paragraphs[0].font.color.rgb = text_color
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    for para in tf.paragraphs:
        para.space_before = Pt(4)
        para.space_after = Pt(4)
    return shape


def add_table_simple(slide, headers, rows, left, top, width, row_height,
                     font_size=11):
    """Add a styled table matching the Excel export aesthetic."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                       Inches(left), Inches(top),
                                       Inches(width),
                                       Inches(row_height * n_rows))
    tbl = tbl_shape.table

    # Header row
    for j, hdr in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = hdr
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.name = FONT_HEAD
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    # Data rows
    for i, row in enumerate(rows):
        fill = ALT_ROW if i % 2 == 1 else SLIDE_BG
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(font_size)
            p.font.name = FONT_BODY
            p.font.color.rgb = TEXT_COLOR
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill

    return tbl_shape


def add_flow_step(slide, num, label, left, top, width=2.6, height=0.7):
    """Single step box in a workflow diagram."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECONDARY_BG
    shape.line.color.rgb = PRIMARY
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Step {num}:  "
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = PRIMARY

    run2 = p.add_run()
    run2.text = label
    run2.font.size = Pt(12)
    run2.font.color.rgb = TEXT_COLOR

    return shape


def add_arrow(slide, left, top, width=0.5, height=0.0):
    """Horizontal right-arrow between flow steps."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    return shape


# ══════════════════════════════════════════════════════════════════════════════
# BUILD SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ───────────────────────────────────────────────────
    add_title_slide(
        prs,
        "🏢  Multifamily Demand Index Tool",
        "A Systematic Approach to Market Screening\n"
        "for Value-Add Multifamily Investment  ·  Quannah Partners",
    )

    # ── Slide 2: What Is This? ───────────────────────────────────────────
    slide = add_content_slide(prs, "What Is the Demand Index App?")

    add_bullet_list(slide, [
        "• Standalone web application (Streamlit) — no coding required",
        "• Ingests CoStar Market Analytics exports (.xlsx) for any set of metros",
        "• Optionally enriches data with U.S. Census demographics (population, income, employment, migration)",
        "• Computes a composite 0–100 Demand Index that ranks markets by apartment demand strength relative to supply",
        "• Outputs a fully formula-driven, auditable Excel workbook — every number is traceable",
        "• Classifies markets into tiers:  High Demand (≥ 67)  ·  Moderate (33–66)  ·  Low Demand (< 33)",
    ], 0.5, 1.5, 8.5, 4.5, font_size=14)

    add_callout_box(slide,
        "💡  Designed for analysts — drag-and-drop CoStar files, click one button, "
        "download a ranked workbook ready for investment committee review.",
        0.5, 6.0, 8.5, 0.7, font_size=12, bold=False)

    # Right-side summary card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(9.5), Inches(1.5),
                                  Inches(3.3), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = SECONDARY_BG
    card.line.color.rgb = BORDER_GRAY
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.word_wrap = True
    items = [
        ("Input", "CoStar .xlsx per metro"),
        ("Enrichment", "Census ACS demographics"),
        ("Output", "0–100 ranked index"),
        ("Artifact", "Formula-driven Excel"),
        ("Audience", "Analysts & Partners"),
    ]
    for i, (label, value) in enumerate(items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.space_before = Pt(12)
        run = p.add_run()
        run.text = f"{label}\n"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.name = FONT_HEAD
        run.font.color.rgb = PRIMARY
        run2 = p.add_run()
        run2.text = value
        run2.font.size = Pt(12)
        run2.font.name = FONT_BODY
        run2.font.color.rgb = TEXT_COLOR

    # ── Slide 3: How to Use It ───────────────────────────────────────────
    slide = add_content_slide(prs, "How to Use It", "Four-step workflow from CoStar export to ranked workbook")

    steps = [
        ("Export", "Download quarterly data\nfrom CoStar Market\nAnalytics as .xlsx"),
        ("Upload", "Drag & drop files into\nthe app; optionally enter\na Census API key"),
        ("Compute", "Click 'Compute Demand\nIndex' — the app ingests,\nenriches, scores & ranks"),
        ("Download", "Download the Excel\nworkbook with rankings,\ncomponents & demographics"),
    ]

    x_start = 0.5
    step_w = 2.8
    arrow_w = 0.5
    gap = 0.15
    y_top = 2.2

    for i, (title, desc) in enumerate(steps):
        x = x_start + i * (step_w + arrow_w + gap)
        # Step box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(y_top),
                                       Inches(step_w), Inches(2.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SECONDARY_BG
        shape.line.color.rgb = PRIMARY
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Step {i+1}: {title}"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.name = FONT_HEAD
        run.font.color.rgb = PRIMARY
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = ""
        p2.space_before = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(12)
        p3.font.name = FONT_BODY
        p3.font.color.rgb = TEXT_COLOR
        p3.alignment = PP_ALIGN.CENTER

        # Arrow between steps
        if i < len(steps) - 1:
            ax = x + step_w + gap * 0.3
            add_arrow(slide, ax, y_top + 0.8, width=arrow_w * 0.7)

    add_callout_box(slide,
        "📋  CoStar file naming convention:  MarketName_ST.xlsx  (e.g., Austin_TX.xlsx, Phoenix_AZ.xlsx). "
        "The app derives market names from file names automatically.",
        0.5, 5.0, 12.3, 0.65, font_size=11)

    add_text_box(slide,
        "Census API key: free from census.gov — unlocks employment growth, population growth, "
        "and income growth components for the full 6-factor model.",
        0.5, 5.8, 12.3, 0.5, font_size=11, italic=True, color=NOTE_GRAY)

    # ── Slide 4: Methodology — Index Components ─────────────────────────
    slide = add_content_slide(prs, "Methodology: Index Components",
                              "Two model configurations based on data availability")

    # Left: Full model
    add_text_box(slide, "Full Model (with Census)", 0.4, 1.4, 5.5, 0.4,
                 font_size=16, bold=True, color=PRIMARY)

    full_headers = ["Component", "Weight", "Direction"]
    full_rows = [
        ["Net Absorption (% inv.)", "+25%", "↑ Higher = more demand"],
        ["Employment Growth (YoY)",  "+25%", "↑ Growth = demand driver"],
        ["Population Growth (YoY)",  "+15%", "↑ Future demand signal"],
        ["Income Growth (YoY)",      "+10%", "↑ Affordability capacity"],
        ["Prior-Year Vacancy Rate",  "−15%", "↓ Slack suppresses demand"],
        ["Deliveries (% of inv.)",   "−10%", "↓ Supply glut pressure"],
    ]
    add_table_simple(slide, full_headers, full_rows, 0.4, 1.9, 6.0, 0.38,
                     font_size=11)

    # Right: CoStar-only model
    add_text_box(slide, "CoStar-Only Model (no Census key)", 7.0, 1.4, 5.5, 0.4,
                 font_size=16, bold=True, color=PRIMARY)

    costar_headers = ["Component", "Weight", "Direction"]
    costar_rows = [
        ["Net Absorption (% inv.)", "+40%", "↑ Primary demand signal"],
        ["Occupancy Rate",          "+25%", "↑ Current tightness"],
        ["Prior-Year Vacancy Rate",  "−20%", "↓ Largest suppressive factor"],
        ["Deliveries (% of inv.)",   "−15%", "↓ Supply pressure"],
    ]
    add_table_simple(slide, costar_headers, costar_rows, 7.0, 1.9, 5.8, 0.38,
                     font_size=11)

    add_callout_box(slide,
        "⚠️  Rent growth is intentionally excluded from the index — it is the dependent variable "
        "in the NMHC regressions and including it would create circularity. "
        "Under-construction % is also excluded to avoid double-counting with deliveries.",
        0.4, 5.3, 12.5, 0.65, font_size=11,
        bg_color=MOD_YELLOW, text_color=TEXT_COLOR)

    add_text_box(slide,
        "Positive weights reward demand signals; negative weights penalize supply/slack indicators. "
        "Weight magnitudes are proportional to empirical coefficient sizes from the NMHC panel regression.",
        0.4, 6.1, 12.5, 0.5, font_size=11, italic=True, color=NOTE_GRAY)

    # ── Slide 5: Why This Methodology? ───────────────────────────────────
    slide = add_content_slide(prs, "Why This Methodology?",
                              "Empirically grounded, transparent, and auditable")

    # Left column — research basis
    add_text_box(slide, "Research Foundation", 0.4, 1.4, 5.5, 0.4,
                 font_size=16, bold=True, color=PRIMARY)

    add_bullet_list(slide, [
        "• Weights derived from NMHC December 2024 panel regression coefficient magnitudes — "
          "the largest empirical study of U.S. apartment rent determinants",
        "• Employment growth: +19.8 bps per percentage point (largest positive predictor)",
        "• Net absorption: +16.7 bps per percentage point of inventory",
        "• Vacancy rate: −24.5 to −27.7 bps (largest suppressive factor)",
        "• Deliveries: −5.2 to −7.7 bps; supply exceeded absorption in 93 of 99 quarters studied",
        "• Demand-side framework grounded in De Leeuw (1971) and "
          "Polinsky & Ellwood (1979) housing demand models",
    ], 0.4, 1.9, 5.8, 4.0, font_size=12)

    # Right column — design choices
    add_text_box(slide, "Design Choices", 7.0, 1.4, 5.5, 0.4,
                 font_size=16, bold=True, color=PRIMARY)

    add_bullet_list(slide, [
        "• Z-score standardization ensures fair cross-market comparability regardless of differences "
          "in unit magnitude (e.g., absorption units vs. vacancy percentages)",
        "• 0–100 rescaling produces intuitive scores for non-technical stakeholders",
        "• Tier classification (High / Moderate / Low) enables rapid screening of large market sets",
        "• Formula-driven Excel output makes every calculation auditable — "
          "no black box, every weight is traceable to published research",
        "• Two model configurations accommodate varying data availability "
          "without requiring a Census API key for basic analysis",
    ], 7.0, 1.9, 5.8, 4.0, font_size=12)

    add_callout_box(slide,
        "🔬  Every weight is traceable to published research — not calibrated to our data, "
        "not a black box, not a proprietary score.",
        0.4, 6.0, 12.5, 0.65, font_size=12, bold=True,
        bg_color=HIGH_GREEN, text_color=TEXT_COLOR)

    # ── Slide 6: Scoring Mechanics ───────────────────────────────────────
    slide = add_content_slide(prs, "Scoring Mechanics",
                              "Step-by-step: from raw CoStar data to a 0–100 ranked index")

    calc_steps = [
        ("1. Extract latest quarter",
         "Identify the most recent complete quarter per market, excluding quarter-to-date (QTD) rows."),
        ("2. Compute prior-year vacancy",
         "Look up the vacancy rate from the same quarter one year prior to capture year-over-year supply context."),
        ("3. Cross-sectional z-scores",
         "For each component, compute z = (x − μ) / σ across all uploaded markets. "
         "This standardizes each metric to a common scale."),
        ("4. Apply signed weights",
         "Multiply each z-score by its configured weight (positive for demand signals, negative for supply/slack)."),
        ("5. Sum into raw composite",
         "Add all weighted z-scores to produce a single raw demand score per market."),
        ("6. Rescale to 0–100",
         "Index = (raw − min) / (max − min) × 100, clamping to the [0, 100] range."),
        ("7. Classify into tiers",
         "High Demand ≥ 67   ·   Moderate Demand 33–66   ·   Low Demand < 33"),
    ]

    y = 1.5
    for title, desc in calc_steps:
        # Step title
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(y),
                                      Inches(3.0), Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.name = FONT_HEAD
        p.font.color.rgb = PRIMARY

        # Step description
        tb2 = slide.shapes.add_textbox(Inches(3.7), Inches(y),
                                       Inches(9.0), Inches(0.45))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.name = FONT_BODY
        p2.font.color.rgb = TEXT_COLOR
        y += 0.62

    # Z-score formula callout
    add_callout_box(slide,
        "z = (x − μ) / σ       where μ = mean across markets, σ = standard deviation across markets\n"
        "Index = (raw − min) / (max − min) × 100",
        0.6, 5.9, 7.0, 0.75, font_size=12, bold=False,
        bg_color=SECONDARY_BG, text_color=ACCENT_DARK)

    # Tier color legend
    tier_y = 5.95
    for label, color, tier_x in [
        ("High Demand (≥ 67)", HIGH_GREEN, 8.3),
        ("Moderate (33–66)",   MOD_YELLOW, 10.0),
        ("Low Demand (< 33)",  LOW_RED,    11.6),
    ]:
        swatch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(tier_x), Inches(tier_y),
                                        Inches(1.4), Inches(0.55))
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = color
        swatch.line.color.rgb = BORDER_GRAY
        swatch.line.width = Pt(0.75)
        tf = swatch.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.name = FONT_HEAD
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER

    # ── Slide 7: Limitations & Caveats ───────────────────────────────────
    slide = add_content_slide(prs, "Limitations & Caveats",
                              "Transparent acknowledgment of constraints")

    # Left column — Data limitations
    add_text_box(slide, "Data Limitations", 0.4, 1.4, 5.5, 0.4,
                 font_size=16, bold=True, color=PRIMARY)

    add_bullet_list(slide, [
        "• Census ACS data lags by approximately 1–2 years; the most recent release may "
          "not reflect current conditions",
        "• 5-year estimates smooth over business cycles, potentially masking short-term shifts",
        "• CBSA auto-matching may fail for unusual market name spellings — "
          "manual override is available in the app",
        "• CoStar data quality and coverage vary by market size; "
          "smaller markets may have less reliable inputs",
        "• Prior-year vacancy falls back to current vacancy when insufficient "
          "historical data is available",
    ], 0.4, 1.9, 5.8, 4.0, font_size=12)

    # Right column — Methodological limitations
    add_text_box(slide, "Methodological Limitations", 7.0, 1.4, 5.5, 0.4,
                 font_size=16, bold=True, color=PRIMARY)

    add_bullet_list(slide, [
        "• The index is relative — adding or removing markets changes all scores "
          "since z-scores are computed across the uploaded set",
        "• Weights are static (derived from published research, not dynamically "
          "calibrated to our portfolio data)",
        "• No macro variables included (interest rates, GDP, national unemployment)",
        "• Does not incorporate submarket-level variation — scores are at the MSA level",
        "• Not a price prediction model — measures demand conditions, "
          "not expected rent growth or returns",
    ], 7.0, 1.9, 5.8, 4.0, font_size=12)

    add_callout_box(slide,
        "These are known, bounded limitations. The tool is designed for screening and relative "
        "ranking — it narrows the aperture for deeper underwriting, it does not replace it.",
        0.4, 6.0, 12.5, 0.6, font_size=12, bold=False,
        bg_color=ACCENT_FILL, text_color=ACCENT_DARK)

    # ── Slide 8: Going Forward ───────────────────────────────────────────
    slide = add_content_slide(prs, "How This Helps Going Forward",
                              "From gut feel to data-driven conviction")

    add_bullet_list(slide, [
        "• Enables rapid, repeatable screening of 30+ metros in minutes — "
          "any analyst can run it independently",
        "• Creates an auditable, formula-driven artifact for investment committee review — "
          "no hidden assumptions",
        "• Pairs with the broader ML rent-growth prediction model for two complementary lenses: "
          "fundamental demand (this tool) vs. forward-looking predictive analytics",
        "• Establishes a standardized, research-backed framework that can be extended with "
          "additional data sources (permits, migration, employer announcements, RealPage, Yardi)",
        "• Reduces time-to-conviction when evaluating unfamiliar markets — "
          "move quickly on emerging opportunities before the competition",
        "• Transparent methodology builds trust with LPs and co-investors "
          "who expect quantitative rigor",
    ], 0.5, 1.5, 7.5, 4.5, font_size=14, spacing=10)

    # Right side — value-prop card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(9.0), Inches(1.5),
                                  Inches(3.8), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = SECONDARY_BG
    card.line.color.rgb = PRIMARY
    card.line.width = Pt(1.5)

    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key Value Proposition"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.name = FONT_HEAD
    run.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

    value_items = [
        "Repeatable process",
        "Research-backed weights",
        "Auditable Excel output",
        "No coding required",
        "Expandable framework",
    ]
    for item in value_items:
        p = tf.add_paragraph()
        p.text = f"✓  {item}"
        p.font.size = Pt(13)
        p.font.name = FONT_BODY
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
        p.alignment = PP_ALIGN.LEFT

    add_callout_box(slide,
        "\"From gut feel to data-driven conviction.\"",
        9.0, 5.3, 3.8, 0.6, font_size=14, bold=True,
        bg_color=HIGH_GREEN, text_color=ACCENT_DARK)

    add_text_box(slide,
        "Next steps:  integrate permit pipeline data  ·  "
        "add submarket drill-down  ·  automate quarterly refresh",
        0.5, 6.0, 8.0, 0.6, font_size=12, italic=True, color=NOTE_GRAY)

    # ── Slide 9: Validation — CoStar Backtest Results ────────────────────
    slide = add_content_slide(prs, "Validation: Historical Backtest",
                              "CoStar fundamentals predict forward rent growth (2000–2025)")

    # Left — scatter chart image
    scatter_path = os.path.join(_SCRIPT_DIR, "figures", "validation_scatter_4Q.png")
    if os.path.exists(scatter_path):
        slide.shapes.add_picture(scatter_path,
                                 Inches(0.3), Inches(1.3),
                                 Inches(6.2), Inches(4.2))
    else:
        add_text_box(slide, "[scatter chart image not found]",
                     0.5, 2.5, 5.5, 0.5, font_size=12, italic=True,
                     color=NOTE_GRAY)

    # Right — tier box plot image
    boxplot_path = os.path.join(_SCRIPT_DIR, "figures", "validation_boxplot_4Q.png")
    if os.path.exists(boxplot_path):
        slide.shapes.add_picture(boxplot_path,
                                 Inches(6.7), Inches(1.3),
                                 Inches(6.2), Inches(4.2))
    else:
        add_text_box(slide, "[box plot image not found]",
                     7.0, 2.5, 5.5, 0.5, font_size=12, italic=True,
                     color=NOTE_GRAY)

    # Key stats callout
    add_callout_box(slide,
        "Pearson r = 0.141  ·  p = 2.49 × 10⁻¹⁴  ·  R² = 0.020  ·  "
        "n = 2,900 market-quarters  ·  29 markets  ·  ~100 quarters",
        0.3, 5.7, 8.5, 0.55, font_size=12, bold=True,
        bg_color=HIGH_GREEN, text_color=ACCENT_DARK)

    # Tier summary mini-table
    add_table_simple(slide,
        ["Tier", "n", "Median Fwd 4Q Growth", "Mean", "IQR"],
        [
            ["High Demand",     "748",   "2.18%", "2.49%", "1.04% – 3.53%"],
            ["Moderate Demand", "1,536", "1.87%", "1.87%", "0.46% – 3.08%"],
            ["Low Demand",      "616",   "1.40%", "1.49%", "−0.10% – 2.65%"],
        ],
        0.3, 6.35, 8.5, 0.22, font_size=9)

    # Interpretation note
    add_text_box(slide,
        "Monotonic tier separation: High > Moderate > Low median forward rent growth. "
        "Statistically significant at p < 0.001.",
        9.0, 5.7, 3.8, 1.0, font_size=11, italic=True, color=NOTE_GRAY)

    # ── Slide 10: What Census Data Adds ──────────────────────────────────
    slide = add_content_slide(prs, "The Role of Census Demographics",
                              "Structural demand drivers complement real-time CoStar signals")

    # Left column — Why include Census
    add_text_box(slide, "Why Include Demographics?", 0.5, 1.4, 5.5, 0.4,
                 font_size=16, bold=True, color=PRIMARY)

    add_bullet_list(slide, [
        "Population growth and net migration reveal long-run housing "
          "demand that precedes absorption data",
        "Median household income sets the rent ceiling — "
          "high demand is unsustainable in low-income metros",
        "Census data captures structural trends (education, age, "
          "household formation) invisible to CoStar",
        "Investors and LPs expect macro-demographic due diligence "
          "alongside real-time market fundamentals",
    ], 0.5, 1.9, 5.5, 3.5, font_size=12)

    # Right column — Design insight
    add_text_box(slide, "Validation Insight", 7.0, 1.4, 5.5, 0.4,
                 font_size=16, bold=True, color=PRIMARY)

    add_bullet_list(slide, [
        "CoStar-only fundamentals (absorption, occupancy, deliveries, "
          "vacancy) are the primary drivers of short-term rent growth",
        "Census ACS estimates update annually with a 1–2 year lag, "
          "so they do not improve quarter-ahead prediction",
        "Census variables add strategic context — they help identify "
          "metros with favorable long-run tailwinds vs. headwinds",
        "The full model blends both: CoStar for timing, "
          "Census for conviction on structural demand",
    ], 7.0, 1.9, 5.5, 3.5, font_size=12)

    add_callout_box(slide,
        "Think of CoStar data as the speedometer (how fast is demand moving now?) and Census data "
        "as the fuel gauge (how much runway does this market have?).  Both matter for investment decisions.",
        0.4, 5.6, 12.5, 0.7, font_size=13, bold=False,
        bg_color=ACCENT_FILL, text_color=ACCENT_DARK)

    # ──────────────────────────────────────────────────────────────────────
    output_path = "demand_index_presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved to {output_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build_presentation()
